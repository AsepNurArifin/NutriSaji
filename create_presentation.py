import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # Set slide dimensions to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors matching the design system
    c_primary = RGBColor(199, 91, 57)     # Terracotta (#C75B39)
    c_secondary = RGBColor(44, 62, 80)    # Dark Slate (#2C3E50)
    c_accent = RGBColor(212, 168, 67)     # Gold (#D4A843)
    c_bg_warm = RGBColor(253, 248, 240)   # Cream (#FDF8F0)
    c_bg_pure = RGBColor(255, 255, 255)   # White (#FFFFFF)
    c_text_dark = RGBColor(26, 26, 26)    # Dark Grey (#1A1A1A)
    c_text_light = RGBColor(255, 255, 255) # White
    c_muted = RGBColor(115, 115, 115)     # Gray
    c_red = RGBColor(220, 38, 38)         # Red
    
    blank_layout = prs.slide_layouts[6]
    
    def apply_bg(slide, color):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.fill.background()
        return rect

    def add_title(slide, text, color=c_secondary, size=Pt(36), top=Inches(0.6)):
        txBox = slide.shapes.add_textbox(Inches(0.8), top, Inches(11.7), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Playfair Display"
        p.font.size = size
        p.font.bold = True
        p.font.color.rgb = color
        return txBox

    # SLIDE 1: Halaman Judul
    s1 = prs.slides.add_slide(blank_layout)
    apply_bg(s1, c_bg_warm)
    
    # Decorative terracotta block on the left
    left_block = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    left_block.fill.solid()
    left_block.fill.fore_color.rgb = c_primary
    left_block.line.fill.background()
    
    # Title box
    t_box = s1.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10), Inches(3.5))
    tf1 = t_box.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "NUTRI SAJI"
    p1.font.name = "Playfair Display"
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = c_primary
    
    p2 = tf1.add_paragraph()
    p2.text = "Ramen Tori Paitan 100% Halal di Palembang"
    p2.font.name = "Inter"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = c_secondary
    p2.space_before = Pt(12)
    
    p3 = tf1.add_paragraph()
    p3.text = "Proposal Bisnis Kemitraan Syariah: Zero Gharar, Simetri Informasi, & Tanpa Riba"
    p3.font.name = "Inter"
    p3.font.size = Pt(18)
    p3.font.color.rgb = c_muted
    p3.space_before = Pt(20)

    # SLIDE 2: Latar Belakang & Masalah Pasar (The Problem)
    s2 = prs.slides.add_slide(blank_layout)
    apply_bg(s2, c_bg_pure)
    add_title(s2, "Latar Belakang: Keraguan Halal di Pasar Ramen")
    
    # Left Box - The Problem
    box_left = s2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0))
    tf = box_left.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Mengapa Konsumen Ragu?"
    p.font.name = "Inter"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = c_red
    p.space_after = Pt(14)
    
    points = [
        "Penggunaan Bahan Syubhat/Haram: Restoran ramen Jepang konvensional sangat bergantung pada mirin, sake (alkohol), lard (lemak babi), dan kecap asin non-halal.",
        "Ketiadaan Transparansi Bahan: Kedai menyembunyikan resep dapur, membuat konsumen muslim khawatir akan kontaminasi silang atau bahan tak dikenal.",
        "Harga Abu-Abu (Gharar): Biaya tak terduga seperti service charge dan tax siluman yang baru ditambahkan saat pembayaran di kasir."
    ]
    for pt in points:
        p = tf.add_paragraph()
        p.text = "• " + pt.split(":")[0] + ":"
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = c_text_dark
        p.space_before = Pt(10)
        
        p_desc = tf.add_paragraph()
        p_desc.text = pt.split(":")[1].strip()
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = c_muted
        p_desc.space_after = Pt(10)
        
    # Right Box - Callout
    callout_rect = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.5))
    callout_rect.fill.solid()
    callout_rect.fill.fore_color.rgb = c_bg_warm
    callout_rect.line.color.rgb = c_accent
    callout_rect.line.width = Pt(1.5)
    
    callout_tf = callout_rect.text_frame
    callout_tf.word_wrap = True
    callout_tf.margin_left = callout_tf.margin_top = callout_tf.margin_right = callout_tf.margin_bottom = Inches(0.4)
    
    cp1 = callout_tf.paragraphs[0]
    cp1.text = "DAMPAK BAGI KONSUMEN"
    cp1.font.bold = True
    cp1.font.size = Pt(12)
    cp1.font.color.rgb = c_accent
    
    cp2 = callout_tf.add_paragraph()
    cp2.text = "Adanya rasa cemas (syubhat) saat makan ramen menurunkan kualitas pengalaman kuliner pelanggan muslim. Hal ini membuka peluang bisnis besar bagi kedai yang berani menjamin 100% kehalalan secara radikal dan jujur."
    cp2.font.size = Pt(16)
    cp2.font.color.rgb = c_secondary
    cp2.space_before = Pt(14)

    # SLIDE 3: Solusi Nutri Saji (The Solution)
    s3 = prs.slides.add_slide(blank_layout)
    apply_bg(s3, c_bg_warm)
    add_title(s3, "Solusi Nutri Saji: Rekonstruksi Bahan Baku")
    
    s3_box = s3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    s3_tf = s3_box.text_frame
    s3_tf.word_wrap = True
    
    p = s3_tf.paragraphs[0]
    p.text = "Menghadirkan Cita Rasa Autentik Tanpa Kompromi Syariah"
    p.font.name = "Inter"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = c_primary
    p.space_after = Pt(20)
    
    s3_points = [
        ("Kaldu Paitan Alami", "Kami merebus tulang ayam segar pilihan (bersertifikat syar'i) selama lebih dari 8 jam untuk mengekstrak kolagen alami sehingga menghasilkan kuah kental gurih murni tanpa pengental sintetis."),
        ("Reduksi Apel & Jamur Premium", "Sebagai pengganti mirin dan sake, kami menggunakan kombinasi jus apel Fuji matang, kaldu jamur shiitake kering, dan madu hutan murni yang dipanaskan perlahan."),
        ("Minyak Schmaltz Khusus", "Lemak babi (lard) digantikan dengan lemak kulit ayam yang dicairkan perlahan pada suhu rendah untuk memberikan rasa gurih yang mewah dan bersih."),
        ("Shoyu Halal MUI Asli", "Menggunakan kecap asin Jepang bersertifikasi resmi Halal MUI dengan kandungan alkohol 0.0% sejak awal pembuatan.")
    ]
    
    for i, (title, desc) in enumerate(s3_points):
        col = i % 2
        row = i // 2
        left_pos = Inches(0.8 + col * 5.9)
        top_pos = Inches(2.6 + row * 2.1)
        
        card = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, top_pos, Inches(5.5), Inches(1.8))
        card.fill.solid()
        card.fill.fore_color.rgb = c_bg_pure
        card.line.color.rgb = c_primary
        card.line.width = Pt(1)
        
        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.margin_top = ctf.margin_left = ctf.margin_bottom = ctf.margin_right = Inches(0.2)
        
        cp = ctf.paragraphs[0]
        cp.text = f"{i+1}. {title}"
        cp.font.bold = True
        cp.font.size = Pt(15)
        cp.font.color.rgb = c_secondary
        
        cp_d = ctf.add_paragraph()
        cp_d.text = desc
        cp_d.font.size = Pt(11)
        cp_d.font.color.rgb = c_text_dark
        cp_d.space_before = Pt(6)

    # SLIDE 4: Anatomi Transparansi Bahan Baku (Komparasi Radikal)
    s4 = prs.slides.add_slide(blank_layout)
    apply_bg(s4, c_bg_pure)
    add_title(s4, "Anatomi Transparansi Bahan Baku (Komparasi Radikal)")
    
    cols = [Inches(0.8), Inches(4.5), Inches(8.7)]
    widths = [Inches(3.4), Inches(3.9), Inches(3.9)]
    
    headers = ["KATEGORI", "RAMEN JEPANG KONVENSIONAL", "SOLUSI HALAL NUTRI SAJI"]
    for i, h in enumerate(headers):
        box = s4.shapes.add_textbox(cols[i], Inches(1.6), widths[i], Inches(0.5))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = c_secondary if i == 0 else (c_red if i == 1 else c_primary)
        if i > 0:
            p.alignment = PP_ALIGN.CENTER
            
    rows = [
        ("Kaldu Dasar", "Tonkotsu (Tulang Babi) / MSG Sintetis berlebih", "Tori Paitan (Tulang Ayam Syar'i rebus 8 jam)"),
        ("Pemanis & Aroma", "Mirin & Sake (Alkohol fermentasi berasa manis)", "Ekstrak Apel Fuji & Jamur Shiitake"),
        ("Minyak Kuah", "Lard (Minyak dari Lemak Babi mentah)", "Chicken Schmaltz (Minyak Kulit Ayam murni)"),
        ("Kecap Shoyu", "Shoyu Konvensional dengan residu alkohol aktif", "Shoyu bersertifikat 100% Halal MUI"),
        ("Sistem Akad", "Gharar (Harga belum termasuk tax siluman)", "Pajak Restoran (PB1) sudah nett di awal")
    ]
    
    for row_idx, r in enumerate(rows):
        top_pos = Inches(2.3 + row_idx * 0.95)
        if row_idx % 2 == 1:
            bg_bar = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), top_pos, Inches(11.7), Inches(0.8))
            bg_bar.fill.solid()
            bg_bar.fill.fore_color.rgb = c_bg_warm
            bg_bar.line.fill.background()
            
        for col_idx, text in enumerate(r):
            box = s4.shapes.add_textbox(cols[col_idx], top_pos + Inches(0.05), widths[col_idx], Inches(0.7))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(12)
            if col_idx == 0:
                p.font.bold = True
                p.font.color.rgb = c_secondary
            elif col_idx == 1:
                p.font.color.rgb = c_muted
                p.alignment = PP_ALIGN.CENTER
            else:
                p.font.bold = True
                p.font.color.rgb = c_text_dark
                p.alignment = PP_ALIGN.CENTER

    # SLIDE 5: Pilar Kepatuhan Syariah & Kejelasan Akad (Zero Gharar)
    s5 = prs.slides.add_slide(blank_layout)
    apply_bg(s5, c_bg_warm)
    add_title(s5, "Pilar Kepatuhan Syariah & Kejelasan Akad (Zero Gharar)")
    
    # Left column: Simetri Informasi
    left_box = s5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    
    lp1 = ltf.paragraphs[0]
    lp1.text = "1. Simetri Informasi (Harga Pas)"
    lp1.font.bold = True
    lp1.font.size = Pt(20)
    lp1.font.color.rgb = c_primary
    lp1.space_after = Pt(14)
    
    lp_points = [
        "Harga Nett Sejak Awal: Semua harga menu di brosur, website, dan kedai adalah harga final yang harus dibayar konsumen.",
        "Sudah Termasuk PB1: Tidak ada penambahan biaya pajak restoran (PB1) atau service fee di kasir yang membuat pelanggan bingung.",
        "Transparansi Resep: Semua bahan tertulis jelas di menu. Konsumen berhak mengetahui apa yang masuk ke tubuh mereka."
    ]
    for pt in lp_points:
        p = ltf.add_paragraph()
        p.text = "✔ " + pt.split(":")[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = c_secondary
        p.space_before = Pt(8)
        
        pd = ltf.add_paragraph()
        pd.text = pt.split(":")[1].strip()
        pd.font.size = Pt(13)
        pd.font.color.rgb = c_text_dark
        pd.space_after = Pt(8)

    # Right column: Anti-Riba
    right_box = s5.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.5), Inches(5.0))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    
    rp1 = rtf.paragraphs[0]
    rp1.text = "2. Penolakan Fasilitas Utang Riba"
    rp1.font.bold = True
    rp1.font.size = Pt(20)
    rp1.font.color.rgb = c_primary
    rp1.space_after = Pt(14)
    
    rp_points = [
        "Tidak Menerima Paylater / Kartu Kredit: Pembayaran secara tegas dibatasi hanya menggunakan Debit, Tunai, atau E-Wallet langsung.",
        "Komitmen Akad Bersih: Pelanggan secara sadar menyetujui akad pembelian menggunakan dana miliknya sendiri tanpa bunga.",
        "Bara'ah al-Dhimmah (Pelepasan Tanggung Jawab): Menggeser tanggung jawab teologis dari penjual kepada pembeli secara jelas di kasir."
    ]
    for pt in rp_points:
        p = rtf.add_paragraph()
        p.text = "✔ " + pt.split(":")[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = c_secondary
        p.space_before = Pt(8)
        
        pd = rtf.add_paragraph()
        pd.text = pt.split(":")[1].strip()
        pd.font.size = Pt(13)
        pd.font.color.rgb = c_text_dark
        pd.space_after = Pt(8)

    # SLIDE 6: Daftar Menu Utama & Transparansi Harga
    s6 = prs.slides.add_slide(blank_layout)
    apply_bg(s6, c_bg_pure)
    add_title(s6, "Daftar Menu Utama & Transparansi Harga")
    
    menus = [
        ("Signature Tori Paitan", "Rp 35.000", "Best Seller — Kaldu Kolagen", "Mie lurus kecil, kuah kaldu ayam pekat 8 jam, ayam chashu halal, jamur kikurage, daun bawang, nori.", "Kaldu kolagen kental yang penuh dengan umami alami ayam segar."),
        ("Spicy Tantanmen", "Rp 38.000", "Pedas Nendang — Rasa Wijen", "Mie keriting tebal, kuah kaldu ayam pedas, daging ayam cincang berbumbu, pak choy, wijen sangrai, cabai kering.", "Sensasi kuah pedas wijen gurih dengan topping melimpah."),
        ("Shoyu Ramen Classic", "Rp 32.000", "Aman untuk Anak — Kuah Bening", "Mie lurus kecil, kuah kaldu ayam shoyu halal, telur ajitama setengah matang, nori, jagung manis, daun bawang.", "Ramen kuah shoyu bening yang segar, gurih halus, dan ramah anak.")
    ]
    
    for i, (name, price, badge, ingredients, desc) in enumerate(menus):
        left_pos = Inches(0.8 + i * 4.0)
        
        card = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, Inches(1.8), Inches(3.7), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = c_bg_warm
        card.line.color.rgb = c_primary if i == 0 else c_muted
        card.line.width = Pt(1.5 if i == 0 else 1)
        
        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.margin_top = ctf.margin_left = ctf.margin_bottom = ctf.margin_right = Inches(0.25)
        
        p = ctf.paragraphs[0]
        p.text = badge.upper()
        p.font.bold = True
        p.font.size = Pt(10)
        p.font.color.rgb = c_accent
        
        p_name = ctf.add_paragraph()
        p_name.text = name
        p_name.font.bold = True
        p_name.font.size = Pt(18)
        p_name.font.color.rgb = c_secondary
        p_name.space_before = Pt(6)
        
        p_price = ctf.add_paragraph()
        p_price.text = price
        p_price.font.bold = True
        p_price.font.size = Pt(22)
        p_price.font.color.rgb = c_primary
        p_price.space_before = Pt(8)
        
        p_lbl = ctf.add_paragraph()
        p_lbl.text = "Komposisi Bahan:"
        p_lbl.font.bold = True
        p_lbl.font.size = Pt(11)
        p_lbl.font.color.rgb = c_secondary
        p_lbl.space_before = Pt(12)
        
        p_ing = ctf.add_paragraph()
        p_ing.text = ingredients
        p_ing.font.size = Pt(11)
        p_ing.font.color.rgb = c_text_dark
        p_ing.space_before = Pt(4)
        
        p_desc = ctf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(11)
        p_desc.font.italic = True
        p_desc.font.color.rgb = c_muted
        p_desc.space_before = Pt(10)

    # SLIDE 7: Model Operasional: Menghindari Gharar Kualitas
    s7 = prs.slides.add_slide(blank_layout)
    apply_bg(s7, c_bg_warm)
    add_title(s7, "Model Operasional: Menghindari Gharar Kualitas")
    
    s7_left = s7.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.8), Inches(5.0))
    s7_ltf = s7_left.text_frame
    s7_ltf.word_wrap = True
    
    p = s7_ltf.paragraphs[0]
    p.text = "Makan di Tempat & Ambil Sendiri (Self-Pickup)"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = c_secondary
    p.space_after = Pt(14)
    
    s7_points = [
        "Mencegah Mie Mengembang (Gharar Kualitas): Ramen segar sangat sensitif terhadap waktu. Pengantaran ojek online berisiko membuat mie dingin dan kaldu tidak segar saat tiba di tangan pelanggan.",
        "Komitmen Kepuasan Akad Jual Beli: Kita hanya menyerahkan produk dalam kualitas prima langsung di kedai sehingga tidak ada pihak yang dirugikan.",
        "Dapur Terbuka (Open Kitchen): Menghilangkan keraguan konsumen dengan memperlihatkan proses peracikan secara higienis di depan mata pelanggan."
    ]
    for pt in s7_points:
        p = s7_ltf.add_paragraph()
        p.text = "🎯 " + pt.split(":")[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = c_primary
        p.space_before = Pt(8)
        
        pd = s7_ltf.add_paragraph()
        pd.text = pt.split(":")[1].strip()
        pd.font.size = Pt(13)
        pd.font.color.rgb = c_text_dark
        pd.space_after = Pt(8)
        
    right_rect = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.5))
    right_rect.fill.solid()
    right_rect.fill.fore_color.rgb = c_bg_pure
    right_rect.line.color.rgb = c_accent
    right_rect.line.width = Pt(1.5)
    
    rr_tf = right_rect.text_frame
    rr_tf.word_wrap = True
    rr_tf.margin_top = rr_tf.margin_left = rr_tf.margin_bottom = rr_tf.margin_right = Inches(0.4)
    
    rp = rr_tf.paragraphs[0]
    rp.text = "ALUR OPERASIONAL KEDAI"
    rp.font.bold = True
    rp.font.size = Pt(12)
    rp.font.color.rgb = c_accent
    
    rp2 = rr_tf.add_paragraph()
    rp2.text = "• Lokasi Stand:\nJl. Sudirman, Palembang (Lokasi Strategis)\n\n• Waktu Operasional:\nSelasa s.d. Minggu (16:00 - 22:00 WIB)\nSenin: Libur untuk pembersihan total kedai\n\n• Pemesanan Pra-Kedatangan:\nPelanggan dapat melakukan pemesanan via WhatsApp sebelum datang agar antrean rapi dan mie disiapkan tepat waktu saat tiba."
    rp2.font.size = Pt(14)
    rp2.font.color.rgb = c_secondary
    rp2.space_before = Pt(14)

    # SLIDE 8: Rencana Pemasaran & Target Audiens
    s8 = prs.slides.add_slide(blank_layout)
    apply_bg(s8, c_bg_pure)
    add_title(s8, "Rencana Pemasaran & Target Audiens")
    
    s8_box = s8.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    s8_tf = s8_box.text_frame
    s8_tf.word_wrap = True
    
    p = s8_tf.paragraphs[0]
    p.text = "Strategi Penetrasi Pasar Kuliner Palembang"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = c_secondary
    p.space_after = Pt(20)
    
    marketing_points = [
        ("Kejujuran Ekstrem (Radical Honesty)", "Kampanye pemasaran fokus pada edukasi bahan syubhat kuliner Jepang umum dan menunjukkan solusi halal yang kami gunakan. Transparansi adalah daya tarik utama."),
        ("Target Audiens Terfokus", "Keluarga muslim kelas menengah, mahasiswa/pemuda pecinta budaya pop Jepang, serta komunitas sadar halal (halal-conscious) di Palembang."),
        ("Pemasaran Digital Lokal", "Optimasi Google Maps lokal, ulasan influencer kuliner Palembang berbasis kejujuran, serta konten edukasi anatomi ramen di Instagram & TikTok."),
        ("Efek Getok Tular (Word-of-Mouth)", "Cita rasa kaldu kolagen ayam kental 8 jam yang sangat khas akan mendorong rekomendasi alami dari mulut ke mulut secara organik.")
    ]
    
    for i, (title, desc) in enumerate(marketing_points):
        col = i % 2
        row = i // 2
        left_pos = Inches(0.8 + col * 5.9)
        top_pos = Inches(2.6 + row * 2.1)
        
        card = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, top_pos, Inches(5.5), Inches(1.8))
        card.fill.solid()
        card.fill.fore_color.rgb = c_bg_warm
        card.line.color.rgb = c_accent
        card.line.width = Pt(1)
        
        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.margin_top = ctf.margin_left = ctf.margin_bottom = ctf.margin_right = Inches(0.2)
        
        cp = ctf.paragraphs[0]
        cp.text = f"📢 {title}"
        cp.font.bold = True
        cp.font.size = Pt(14)
        cp.font.color.rgb = c_secondary
        
        cp_d = ctf.add_paragraph()
        cp_d.text = desc
        cp_d.font.size = Pt(12)
        cp_d.font.color.rgb = c_text_dark
        cp_d.space_before = Pt(6)

    # SLIDE 9: Analisis Keuangan & Proyeksi BEP (Financial Projection)
    s9 = prs.slides.add_slide(blank_layout)
    apply_bg(s9, c_bg_warm)
    add_title(s9, "Analisis Keuangan & Proyeksi BEP")
    
    stats = [
        ("Investasi Awal", "Rp 85 Juta", "Sewa tempat, renovasi kedai terbuka, peralatan dapur, dan modal kerja awal."),
        ("Margin Kotor", "60% - 65%", "Biaya bahan baku per mangkok rata-rata berkisar di angka Rp 12.000 - Rp 14.000."),
        ("Margin Bersih", "30% - 35%", "Setelah dikurangi biaya sewa, gaji karyawan, utilitas kedai, dan penyusutan alat.")
    ]
    
    for i, (title, val, desc) in enumerate(stats):
        left_pos = Inches(0.8 + i * 4.0)
        
        card = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, Inches(1.8), Inches(3.7), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = c_bg_pure
        card.line.color.rgb = c_primary
        card.line.width = Pt(1.5)
        
        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.margin_top = ctf.margin_left = ctf.margin_bottom = ctf.margin_right = Inches(0.2)
        
        p = ctf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = c_secondary
        p.alignment = PP_ALIGN.CENTER
        
        p_val = ctf.add_paragraph()
        p_val.text = val
        p_val.font.bold = True
        p_val.font.size = Pt(28)
        p_val.font.color.rgb = c_primary
        p_val.space_before = Pt(8)
        p_val.alignment = PP_ALIGN.CENTER
        
        p_desc = ctf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(10)
        p_desc.font.color.rgb = c_muted
        p_desc.space_before = Pt(8)
        p_desc.alignment = PP_ALIGN.CENTER
        
    timeline_box = s9.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.5))
    ttf = timeline_box.text_frame
    ttf.word_wrap = True
    
    tp = ttf.paragraphs[0]
    tp.text = "PROYEKSI BALIK MODAL (BREAK-EVEN POINT)"
    tp.font.bold = True
    tp.font.size = Pt(14)
    tp.font.color.rgb = c_secondary
    tp.space_after = Pt(10)
    
    tp2 = ttf.add_paragraph()
    tp2.text = "• Target Penjualan Harian: Rata-rata 50 mangkok per hari di 3 bulan pertama, meningkat ke 80 mangkok per hari di bulan ke-6.\n" \
               "• Omzet Bulanan Target: Rp 52.500.000 - Rp 84.000.000.\n" \
               "• BEP Operasional Kedai: Bulan ke-2 (arus kas positif sejak awal).\n" \
               "• Balik Modal Investasi (ROI): Ditargetkan penuh pada bulan ke-10 s.d. bulan ke-12 operasional kedai."
    tp2.font.size = Pt(13)
    tp2.font.color.rgb = c_text_dark
    tp2.space_before = Pt(6)

    # SLIDE 10: Visi Masa Depan & Peluang Syirkah (Kemitraan)
    s10 = prs.slides.add_slide(blank_layout)
    apply_bg(s10, c_secondary)
    
    r_block = s10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(0), Inches(5.833), Inches(7.5))
    r_block.fill.solid()
    r_block.fill.fore_color.rgb = c_bg_warm
    r_block.line.fill.background()
    
    l_box = s10.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(6.0), Inches(5.0))
    ltf = l_box.text_frame
    ltf.word_wrap = True
    
    lp = ltf.paragraphs[0]
    lp.text = "Visi Jangka Panjang"
    lp.font.name = "Playfair Display"
    lp.font.size = Pt(36)
    lp.font.bold = True
    lp.font.color.rgb = c_accent
    
    lp2 = ltf.add_paragraph()
    lp2.text = "Menjadi pelopor kedai ramen 100% halal dan transparan yang dipercaya oleh masyarakat Palembang, lalu memperluas keberkahan ini ke kota-kota lain di Indonesia melalui skema kemitraan syirkah bebas riba."
    lp2.font.size = Pt(16)
    lp2.font.color.rgb = c_text_light
    lp2.space_before = Pt(20)
    
    lp3 = ltf.add_paragraph()
    lp3.text = "Kami mengundang Anda bergabung sebagai mitra pemodal dalam skema Mudharabah yang adil dan bersih dari akad bathil."
    lp3.font.size = Pt(14)
    lp3.font.color.rgb = c_accent
    lp3.space_before = Pt(20)

    r_box = s10.shapes.add_textbox(Inches(8.0), Inches(1.5), Inches(4.8), Inches(5.0))
    rtf = r_box.text_frame
    rtf.word_wrap = True
    
    rp = rtf.paragraphs[0]
    rp.text = "Hubungi Kemitraan"
    rp.font.name = "Inter"
    rp.font.size = Pt(24)
    rp.font.bold = True
    rp.font.color.rgb = c_primary
    
    rp2 = rtf.add_paragraph()
    rp2.text = "Mari berdiskusi lebih lanjut untuk menyatukan visi kuliner halal ini."
    rp2.font.size = Pt(13)
    rp2.font.color.rgb = c_muted
    rp2.space_before = Pt(10)
    
    rp3 = rtf.add_paragraph()
    rp3.text = "📍 Lokasi Kedai:\nJl. Sudirman (Samping X), Palembang\n\n💬 Kontak WhatsApp:\n+62 811-1222-333\n\n🌐 Website Catalog:\nwww.nutrisaji.com"
    rp3.font.size = Pt(14)
    rp3.font.bold = True
    rp3.font.color.rgb = c_secondary
    rp3.space_before = Pt(20)
    
    prs.save("Proposal_Bisnis_Nutri_Saji.pptx")
    print("Presentation created successfully as 'Proposal_Bisnis_Nutri_Saji.pptx'!")

if __name__ == "__main__":
    create_deck()
